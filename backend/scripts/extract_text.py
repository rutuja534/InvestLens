# scripts/extract_text.py
import sys
from pptx import Presentation
import pytesseract
from PIL import Image
import io

def extract_text_from_pptx(filepath):
    prs = Presentation(filepath)
    full_text = []

    for slide in prs.slides:
        slide_text = []
        # Extract text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        
        # Extract text from images using OCR
        for shape in slide.shapes:
            if shape.shape_type == 13: # 13 corresponds to a picture
                try:
                    image = shape.image
                    image_bytes = image.blob
                    img = Image.open(io.BytesIO(image_bytes))
                    ocr_text = pytesseract.image_to_string(img)
                    slide_text.append(ocr_text)
                except Exception as e:
                    # Silently pass if an image can't be processed
                    pass

        full_text.append('\n'.join(slide_text))

    # Join text from all slides with a separator
    return "\n--- SLIDE BREAK ---\n".join(full_text)

if __name__ == "__main__":
    # The first command-line argument is the script name, the second is our file path
    pptx_file_path = sys.argv[1]
    extracted_text = extract_text_from_pptx(pptx_file_path)
    print(extracted_text)